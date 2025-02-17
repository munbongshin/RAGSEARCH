import sys, os
import io, re
import tempfile
import time
import PyPDF2
import docx
import logging
import pdfplumber
import json

from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import markdown
import html
import pandas as pd
from typing import Union, Optional, List, IO, Dict, Any
from langchain.docstore.document import Document
from dotenv import load_dotenv
from pathlib import Path
import hwp5
from hwp5.xmlmodel import Hwp5File
import platform
import numpy as np

if platform.system() == 'Windows':
    import win32com.client as win32    
    import win32gui
    import win32con    


logging.basicConfig(level=logging.ERROR)
logger = logging.getLogger(__name__)

class TableProcessor:
    def __init__(self):
        """TableProcessor 클래스 초기화"""
        self.logger = logging.getLogger(__name__)

    def process_table(self, table: List[List[str]]) -> str:
        """
        모든 형태의 데이터를 마크다운 표 형식으로 변환하는 범용 처리기
        
        Args:
            table (List[List[str]]): 2차원 리스트 형태의 표 데이터
        Returns:
            str: 마크다운 형식의 표
        """
        try:
            if not table or not any(any(cell for cell in row) for row in table):
                return ""

            processed_table = self._preprocess_data(table)
            col_widths = self._calculate_column_widths(processed_table)
            return self._generate_markdown_table(processed_table, col_widths)

        except Exception as e:
            self.logger.error(f"표 처리 중 오류 발생: {str(e)}")
            return str(table)

    def _preprocess_data(self, table: List[List[str]]) -> List[List[str]]:
        """데이터 전처리"""
        processed = []
        for row in table:
            if not any(cell for cell in row):
                continue
                
            processed_row = []
            for cell in row:
                cell_str = str(cell).strip() if cell else ""
                
                if re.search(r'\d{1,2}:\d{2}', cell_str):
                    if '-' in cell_str:
                        times = cell_str.split('-')
                        if len(times) == 2:
                            cell_str = f"{times[0].strip()}-{times[1].strip()}"
                
                if '(' in cell_str and ')' in cell_str:
                    main_text = cell_str.split('(')[0].strip()
                    info = cell_str[cell_str.find('('):]
                    cell_str = f"{main_text}<br>{info}"
                
                processed_row.append(cell_str)
            processed.append(processed_row)
        
        return processed

    def _calculate_column_widths(self, table: List[List[str]]) -> List[int]:
        """각 열의 최대 너비 계산"""
        if not table:
            return []
            
        widths = []
        for col_idx in range(max(len(row) for row in table)):
            col_cells = [row[col_idx] if col_idx < len(row) else "" for row in table]
            width = max((len(str(cell).split('<br>')[0]) for cell in col_cells), default=0)
            widths.append(width)
            
        return widths

    def _generate_markdown_table(self, table: List[List[str]], col_widths: List[int]) -> str:
        """마크다운 표 생성"""
        if not table:
            return ""

        header = table[0]
        md_table = [
            "| " + " | ".join(str(cell).ljust(width) for cell, width in zip(header, col_widths)) + " |",
            "|" + "|".join("-" * (width + 2) for width in col_widths) + "|"
        ]

        for row in table[1:]:
            formatted_cells = []
            for col_idx, cell in enumerate(row):
                if col_idx < len(col_widths):
                    if '<br>' in str(cell):
                        main_text, info = str(cell).split('<br>')
                        formatted_cell = f"{main_text.ljust(col_widths[col_idx])}<br>{info}"
                    else:
                        formatted_cell = str(cell).ljust(col_widths[col_idx])
                    formatted_cells.append(formatted_cell)
            
            md_table.append("| " + " | ".join(formatted_cells) + " |")

        return "\n".join(md_table)

class ExtractTextFromFile:
    def __init__(self):
        project_root = Path(__file__).parent.parent
        # .env 파일의 경로 설정
        env_path = project_root / '.env'
        load_dotenv(dotenv_path=env_path)
        self.chunk_size = os.environ.get("CHUNK_SIZE")
        self.chunk_overlap = os.environ.get("CHUNK_OVERLAP")
        self.table = TableProcessor()  # TableProcessor 인스턴스 생성
    
    def process_table_data(self, data: List[List[str]]) -> str:
        return self.table.process_table(data)

    
    def extract_text_from_pdf_pages(self, file: Union[str, IO],file_name) -> List[Document]:
        """
        Extract text from PDF pages.

        Args:
            file (Union[str, IO]): Path to the PDF file or a file-like object

        Returns:
            List[Document]: List of Document objects containing extracted text from PDF pages
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        # 파일 객체가 읽기 모드로 열려 있는지 확인
        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        file_obj = None
        
        try:
            if isinstance(file, str):
                file_path = file
                file_name = file_name
                file_obj = open(file, 'rb')
            else:
                file_obj = file
                file_name = file_name
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우

            pdf_reader = PyPDF2.PdfReader(file_obj)

            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                text = self.clean_text2(text)                
                metadata = {
                    "source": os.path.basename(file_name),
                    "file_name": file_name,
                    "page": page_num
                }
                doc = Document(page_content=text, source=os.path.basename(file_name), metadata=metadata)
                documents.append(doc)

            return documents

        except Exception as e:
            logger.debug(f"Error processing PDF file: {str(e)}")
            return []

        finally:
            # 우리가 파일을 열었다면 닫아줍니다
            if file_obj and isinstance(file_obj, io.IOBase) and not isinstance(file_obj, io.BytesIO):
                file_obj.close()
                 
    def extract_text_from_pdf_pages_plumber(self, file: Union[str, IO], file_name: str) -> List[Document]:
        """
        PDF에서 텍스트 추출
        """
        documents = []
        file_obj = None
        
        try:
            if isinstance(file, str):
                file_path = file_name
                file_obj = open(file, 'rb')
            else:
                file_obj = file
                file_path = file_name

            with pdfplumber.open(file_obj) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # 일반 텍스트 추출
                    text = page.extract_text() or ""
                    
                    # 표 추출 및 처리
                    tables = page.extract_tables()
                    for table in tables:
                        if table and any(any(cell for cell in row) for row in table):
                            table_text = self.process_table_data(table)
                            if table_text:
                                text += f"\n{table_text}\n"
                    
                    # 텍스트 정리
                    text = self.clean_text2(text)
                    
                    # 메타데이터 설정
                    metadata = {
                        "source": os.path.basename(file_name),
                        "file_name": file_name,
                        "page": page_num,
                        "tables_found": len(tables)
                    }
                    
                    doc = Document(page_content=text, metadata=metadata)
                    documents.append(doc)
                    
            return documents

        except Exception as e:
            logger.error(f"PDF 처리 중 오류 발생: {str(e)}")            
            return []

        finally:
            if file_obj and isinstance(file_obj, io.IOBase) and not isinstance(file_obj, io.BytesIO):
                file_obj.close()

    def _analyze_table_structure(self, table: List[List[str]]) -> dict:
        """
        테이블 구조 분석
        
        Args:
            table (List[List[str]]): 분석할 테이블 데이터
            
        Returns:
            dict: 테이블 구조 정보
        """
        structure = {
            'has_merged_header': False,
            'header_rows': [],
            'main_groups': set(),
            'sub_groups': set()
        }
        
        if not table:
            return structure

        # 첫 번째 열 분석 (메인 그룹)
        first_col = [row[0] if row else '' for row in table]
        unique_first_col = set(cell.strip() for cell in first_col if cell and cell.strip())
        structure['main_groups'] = unique_first_col

        # 병합된 헤더 패턴 확인
        for row_idx, row in enumerate(table):
            # 같은 값이 반복되는지 확인
            values = {}
            for cell in row:
                if cell and str(cell).strip():
                    cell_value = str(cell).strip()
                    values[cell_value] = values.get(cell_value, 0) + 1
            
            # 반복되는 값이 있으면 병합된 헤더로 간주
            if any(count > 1 for count in values.values()):
                structure['has_merged_header'] = True
                structure['header_rows'].append(row_idx)

        return structure
    
    

    def _process_merged_header_table(self, table: List[List[str]]) -> str:
        """
        병합된 헤더가 있는 테이블 처리
        
        Args:
            table (List[List[str]]): 처리할 테이블 데이터
            
        Returns:
            str: 처리된 테이블 텍스트
        """
        table_text = []
        current_main_group = None
        
        for row_idx, row in enumerate(table):
            row_cells = [str(cell).strip() if cell else "" for cell in row]
            
            # 빈 행 건너뛰기
            if not any(row_cells):
                continue

            # 메인 그룹 처리
            if row_cells[0] and row_cells[0] != current_main_group:
                current_main_group = row_cells[0]
                table_text.append(f"\n{current_main_group}")

            # 데이터 행 처리
            row_text = []
            for col_idx, cell in enumerate(row_cells[1:], 1):
                if not cell:
                    continue
                    
                # 시간 정보 특별 처리
                if re.search(r'\d{1,2}[:.]\d{2}', cell):
                    row_text.append(cell)
                else:
                    row_text.append(cell)

            if row_text:
                if current_main_group:
                    table_text.append("  " + ", ".join(row_text))
                else:
                    table_text.append(", ".join(row_text))

        return "\n".join(filter(None, table_text))

    def _process_standard_table(self, table: List[List[str]]) -> str:
        """
        일반적인 표 처리
        
        Args:
            table (List[List[str]]): 처리할 테이블 데이터
            
        Returns:
            str: 처리된 테이블 텍스트
        """
        table_text = []
        headers = [str(cell).strip() if cell else "" for cell in table[0]]
        
        for row_idx, row in enumerate(table[1:], 1):
            row_cells = [str(cell).strip() if cell else "" for cell in row]
            if not any(row_cells):
                continue

            row_text = []
            for col_idx, cell in enumerate(row_cells):
                if not cell:
                    continue
                    
                if col_idx < len(headers) and headers[col_idx]:
                    row_text.append(f"{headers[col_idx]}: {cell}")
                else:
                    row_text.append(cell)

            if row_text:
                table_text.append(" | ".join(row_text))

        return "\n".join(table_text)
    
    def clean_text2(self, text) -> str:
        try:
            # 1. bytes 처리를 가장 먼저
            if isinstance(text, bytes):
                try:
                    text = text.decode('utf-8')
                except UnicodeDecodeError:
                    try:
                        text = text.decode('cp949')
                    except UnicodeDecodeError:
                        logger.warning(f"Failed to decode bytes: {text}")
                        return str(text)

            # 원본이 None이거나 빈 문자열인 경우
            if not text:
                return ""
                    
            # 문자열로 변환
            text = str(text).strip()
            
            # 중복된 특수 문자 제거
            text = re.sub(r'\|{2,}', '|', text)  # 연속된 | 기호 하나로 줄이기
            text = re.sub(r'-{2,}', ' ', text)   # 연속된 - 기호 공백으로 대체

            # 마크다운 형식으로 전처리
            text = re.sub(r'\n{3,}', '\n\n', text)  # 연속된 줄바꿈 2줄까지만 유지
            text = re.sub(r'^(\d+)\.\s*(.+)$', r'**\1. \2**', text, flags=re.MULTILINE)  # 번호 목록 강조
            text = re.sub(r'^(.+):$', r'**\1:**', text, flags=re.MULTILINE)  # 콜론 라인 강조

            # HTML 이스케이프 처리
            text = html.escape(text)
            html_content=text
            # Markdown으로 변환
            """html_content = markdown.markdown(
                text, 
                extensions=[
                    'markdown.extensions.fenced_code',
                    'markdown.extensions.tables',
                    'markdown.extensions.nl2br'  # 줄바꿈 처리
                ],
                output_format='html5'
            )"""

            return html_content.strip()
                
        except Exception as e:
            logger.error(f"텍스트 정제 중 오류 발생: {str(e)}")
            return str(text)

    def clean_text(self, text) -> str:
        """
        텍스트 정제 함수
        
        Args:
            text: 정제할 텍스트 (str 또는 bytes)
            
        Returns:
            str: 정제된 텍스트
        """
        try:
            # 1. bytes 처리를 가장 먼저
            if isinstance(text, bytes):
                try:
                    text = text.decode('utf-8')
                except UnicodeDecodeError:
                    try:
                        text = text.decode('cp949')
                    except UnicodeDecodeError:
                        logger.warning(f"Failed to decode bytes: {text}")
                        return str(text)

            # 원본이 None이거나 빈 문자열인 경우
            if not text:
                return ""
                
            # 문자열로 변환
            text = str(text).strip()
            
            # 새로 추가된 부분: 연속된 점, |, ○ 처리
            text = re.sub(r'\.{2,}', '.', text)  # 연속된 점 하나로 줄이기
            text = re.sub(r'\|{2,}', '|', text)  # 연속된 | 하나로 줄이기
            text = re.sub(r'○', ' ', text)       # ○ 문자 공백으로 대체
            text = re.sub(r'[|.]{2,}', ' ', text)  # 연속된 특수문자 공백으로 대체
            
            # 2. 여러 개의 공백을 하나로 통합 (줄바꿈 제외)
            text = re.sub(r'[^\S\n]+', ' ', text)
            
            # 3. 시간 형식의 콜론 주변 공백 처리
            text = re.sub(r'(\d)\s*:\s*(\d)', r'\1:\2', text)
            
            # 4. 숫자 사이의 물결 표시 주변 공백 제거
            text = re.sub(r'(\d)\s*~\s*(\d)', r'\1~\2', text)
            
            # 5. 불필요한 유니코드 문자 제거
            text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)

            # 6. | 문자를 공백으로 대체 (앞뒤에 글자가 있는 경우만 공백 유지)
            text = re.sub(r'(\w)\s*\|\s*(\w)', r'\1 \2', text)  # 글자 사이의 | 
            text = re.sub(r'^\|\s*', '', text)                  # 문장 시작의 |
            text = re.sub(r'\s*\|$', '', text)                  # 문장 끝의 |
            
            # 7. 하이픈(-) 문자가 단독으로 있을 때만 공백으로 대체
            text = re.sub(r'(\w)\s*-\s*(\w)', r'\1 \2', text)  # 글자 사이의 -
            text = re.sub(r'^\-\s*', '', text)                 # 문장 시작의 -
            text = re.sub(r'\s*\-$', '', text)                 # 문장 끝의 -

            # 8. □, ○, ●, ◦ 문자를 공백으로 대체
            text = re.sub(r'(\w)[□○●◦]\s*(\w)', r'\1 \2', text)  # 글자 사이의 특수문자
            text = re.sub(r'^[□○●◦]\s*', '', text)              # 문장 시작의 특수문자
            text = re.sub(r'\s*[□○●◦]$', '', text)              # 문장 끝의 특수문자
            
            # 9. 여러 개의 공백을 하나로 통합 (줄바꿈 제외)
            text = re.sub(r'[^\S\n]+', ' ', text)
            text = re.sub(r'-{2,}', ' ', text) 
            
            # 10. 문자열 앞뒤 공백 제거
            return text.strip()
            
        except Exception as e:
            logger.error(f"텍스트 정제 중 오류 발생: {str(e)}")
            logger.error(f"원본 텍스트: {repr(text)}")
            if isinstance(text, bytes):
                try:
                    return text.decode('utf-8', errors='ignore')
                except:
                    return str(text)
            return str(text)
        

    def extract_text_from_docx(self, file):
        doc = Document(file)
        return "\n".join([para.text for para in doc.paragraphs])
    
    def extract_text_from_docx_pages(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from Word document pages.

        Args:
            file (Union[str, IO]): Path to the Word document or a file-like object

        Returns:
            List[Document]: List of Document objects containing extracted text from Word document pages
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        # 파일 객체가 읽기 모드로 열려 있는지 확인
        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        temp_file = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file)
                doc = docx.Document(file)
            else:
                # 임시 파일 생성
                temp_file = io.BytesIO(file.read())
                file_name = self.get_file_name(file)
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우
                doc = docx.Document(temp_file)

            current_page = ""
            page_number = 1

            for para in doc.paragraphs:
                current_page += para.text + "\n"
                current_page = self.clean_text2(current_page)
                
                # 임의로 1000자마다 새 페이지로 간주 (실제 페이지 구분은 더 복잡할 수 있음)
                if len(current_page) > 1000:
                    metadata = {
                        "source": os.path.basename(file_name),
                        "file_name": file_name,
                        "page": page_number
                    }
                    document = Document(page_content=current_page.strip(), metadata=metadata)
                    documents.append(document)
                    current_page = ""
                    page_number += 1

            # 마지막 페이지 추가
            if current_page:
                metadata = {
                    "source": os.path.basename(file_name),
                    "file_name": file_name,
                    "page": page_number
                }
                document = Document(page_content=current_page.strip(), metadata=metadata)
                documents.append(document)

            return documents

        except Exception as e:
            logger.debug(f"Error processing Word document: {str(e)}")
            return []

        finally:
            if temp_file:
                temp_file.close()   
                
    def extract_text_from_xlsx(self, file_path):
        wb = load_workbook(file_path)
        text = ""
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"                
        return text
    
    def extract_text_from_xlsx_pages(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from Excel file pages.

        Args:
            file (Union[str, IO]): Path to the Excel file or a file-like object

        Returns:
            List[Document]: List of Document objects containing extracted text from Excel sheets
        """
        documents = []
        temp_file = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file_path)
            else:
                # 임시 파일 생성
                temp_file = io.BytesIO(file.read())
                file_name = self.get_file_name(file)
                file_path = file_name

            # ExcelFile 객체를 사용하여 모든 시트를 한 번에 읽습니다
            with pd.ExcelFile(temp_file if temp_file else file_path) as xls:
                for sheet_name in xls.sheet_names:
                    # 각 시트를 DataFrame으로 읽습니다
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    
                    if df.empty:
                        logger.warning(f"Empty sheet found: {sheet_name}")
                        continue
                    
                    # DataFrame을 분석 가능한 텍스트로 변환합니다
                    text_chunks = self._convert_df_to_text(df)
                    # 리스트를 문자열로 변환
                    sheet_content = " ".join(text_chunks) if isinstance(text_chunks, list) else text_chunks

                    # Document 객체를 생성합니다
                    metadata = {
                        "source": os.path.basename(file_name),
                        "file_name": file_name,
                        "page": sheet_name
                    }
                    doc = Document(page_content=sheet_content, metadata=metadata)
                    documents.append(doc)

            if not documents:
                logger.warning(f"No content extracted from file: {file_name}")
                
            return documents

        except Exception as e:
            logger.error(f"Error processing Excel file: {str(e)}")
            return []

        finally:
            if temp_file:
                temp_file.close()
    
    
    def analyze_header_importance(self, df: pd.DataFrame) -> Dict[str, float]:
        """
        데이터 분석을 통해 각 헤더(컬럼)의 중요도를 자동으로 계산합니다.
        
        Args:
            df (pd.DataFrame): 분석할 데이터프레임
            
        Returns:
            Dict[str, float]: 컬럼별 중요도 점수 (0~1 사이 정규화된 값)
        """
        try:
            if df.empty:
                logger.warning("Empty DataFrame provided")
                return {}
                
            importance_scores = {}
            
            for column in df.columns:
                try:
                    col_data = df[column].dropna()
                    score = 0.0
                    
                    # 1. 고유성 분석 (최대 2점)
                    if len(col_data) > 0:
                        unique_ratio = len(col_data.unique()) / len(col_data)
                        score += unique_ratio * (2 if 0.3 <= unique_ratio <= 0.9 else 1)
                    
                    # 2. 데이터 완전성 (최대 1점)
                    completion_ratio = 1 - (df[column].isna().sum() / len(df))
                    score += completion_ratio
                    
                    # 3. 데이터 패턴 분석 (최대 2점)
                    sample_values = col_data.head(min(10, len(col_data))).astype(str)
                    
                    # 3.1 길이 일관성 (최대 1점)
                    if not sample_values.empty:
                        lengths = [len(str(v)) for v in sample_values]
                        length_variance = np.var(lengths) if len(lengths) > 1 else 0
                        pattern_consistency = 1 / (1 + length_variance)
                    else:
                        pattern_consistency = 0
                    
                    # 3.2 형식 일관성 (최대 1점)
                    def get_value_pattern(value):
                        if pd.isna(value):
                            return 'null'
                        value = str(value).strip()
                        if not value:
                            return 'empty'
                        if value.isdigit():
                            return 'numeric'
                        if value.replace('.', '', 1).isdigit():  # 소수점 하나만 허용
                            return 'decimal'
                        if any(c.isdigit() for c in value) and any(c.isalpha() for c in value):
                            return 'alphanumeric'
                        if all(c.isalpha() or c.isspace() for c in value):
                            return 'text'
                        return 'mixed'
                    
                    if not sample_values.empty:
                        patterns = [get_value_pattern(v) for v in sample_values]
                        pattern_ratio = len(set(patterns)) / len(patterns) if patterns else 1
                        pattern_consistency += (1 - pattern_ratio)
                    
                    score += pattern_consistency
                    
                    # 4. 참조 빈도 분석 (최대 1점)
                    if not col_data.empty:
                        reference_score = 0
                        sample_size = min(5, len(col_data))
                        sample_values = col_data.head(sample_size)
                        
                        for other_col in df.columns:
                            if other_col != column:
                                try:
                                    # 문자열 비교 전 데이터 타입 통일
                                    other_col_data = df[other_col].astype(str)
                                    ref_count = sum(
                                        other_col_data.str.contains(
                                            str(val), regex=False, na=False
                                        ).sum()
                                        for val in sample_values
                                    )
                                    reference_score += ref_count / (sample_size * len(df))
                                except Exception as e:
                                    logger.warning(f"Reference analysis failed for {column}->{other_col}: {str(e)}")
                                    continue
                        
                        score += min(reference_score, 1.0)  # 최대 1점으로 제한
                    
                    importance_scores[column] = score
                    
                except Exception as e:
                    logger.error(f"Error analyzing column {column}: {str(e)}")
                    importance_scores[column] = 0.0
            
            # 점수 정규화
            if importance_scores:
                max_score = max(importance_scores.values())
                if max_score > 0:
                    importance_scores = {
                        col: score/max_score 
                        for col, score in importance_scores.items()
                    }
                else:
                    # 모든 점수가 0인 경우
                    importance_scores = {col: 1.0/len(importance_scores) for col in importance_scores}
            
            logger.debug(f"Importance scores before normalization: {importance_scores}")
            return importance_scores
            
        except Exception as e:
            logger.error(f"Header importance analysis failed: {str(e)}")
            return {col: 1.0/len(df.columns) for col in df.columns}  # 균등한 가중치 반환

    def detect_key_columns(self, df: pd.DataFrame) -> List[str]:
        """
        데이터 분석을 통해 키 컬럼을 자동으로 감지합니다.
        
        Args:
            df (pd.DataFrame): 분석할 데이터프레임
            
        Returns:
            List[str]: 감지된 키 컬럼 리스트
        """
        try:
            if df.empty:
                logger.warning("Empty DataFrame provided")
                return []
                
            # 컬럼 중요도 점수 계산
            importance_scores = self.analyze_header_importance(df)
            if not importance_scores:
                logger.warning("No importance scores calculated")
                return [df.columns[0]] if len(df.columns) > 0 else []
                
            logger.info(f"Column importance scores: {importance_scores}")
            
            # threshold 계산
            scores = sorted(importance_scores.values(), reverse=True)
            if len(scores) <= 3:
                threshold = 0.5
            else:
                index = max(1, int(len(scores) * 0.3))  # 최소 1개는 선택되도록
                threshold = scores[index - 1]  # 0-based index 조정
            
            logger.debug(f"Using threshold: {threshold}")
            
            # 키 컬럼 선택
            key_columns = [
                col for col, score in importance_scores.items() 
                if score >= threshold
            ]
            
            # 키 컬럼이 없는 경우 가장 높은 점수의 컬럼 선택
            if not key_columns and importance_scores:
                max_col = max(importance_scores.items(), key=lambda x: x[1])[0]
                key_columns = [max_col]
                logger.info(f"No columns above threshold, using highest scoring column: {max_col}")
            
            logger.info(f"Detected {len(key_columns)} key columns: {key_columns}")
            return key_columns
            
        except Exception as e:
            logger.error(f"Error detecting key columns: {str(e)}")
             # 에러 발생 시 첫 번째 컬럼을 반환
            return [df.columns[0]] if len(df.columns) > 0 else []

    def clean_value(self, value: Any) -> str:
        """
        다양한 타입의 값을 텍스트로 변환
        
        Args:
            value: 변환할 값 (Any type)
        Returns:
            str: 정제된 문자열 값
        """
        try:
            # None 또는 NaN 처리
            if value is None or pd.isna(value):
                return ""
                
            # 날짜/시간 처리
            if isinstance(value, (pd.Timestamp, np.datetime64)):
                try:
                    return pd.to_datetime(value).strftime(self.date_format)
                except Exception as e:
                    logger.warning(f"Date conversion error: {e}")
                    return str(value)
                    
            # 숫자 처리
            if isinstance(value, (int, float)):
                if pd.isna(value):  # np.nan 체크
                    return ""
                # 정수인 경우
                if float(value).is_integer():
                    return str(int(value))
                # 소수인 경우
                return f"{value:.2f}"
                
            # 문자열 처리
            str_value = str(value).strip()
            if not str_value:  # 빈 문자열 체크
                return ""
                
            return str_value
            
        except Exception as e:
            logger.error(f"Error in clean_value: {str(e)}, value: {value}, type: {type(value)}")
            return str(value) if value is not None else ""

    def _convert_df_to_text(self, df: pd.DataFrame) -> str:
        chunks = []
        key_columns = self.detect_key_columns(df)
        
        for idx, row in df.iterrows():
            sentences = []
            
            # 키 컬럼 처리
            if key_columns:
                key_info = []
                for key_col in key_columns:
                    value = self.clean_value(row[key_col])
                    key_info.append(f"{key_col} {value}")
                sentences.append(", ".join(key_info))
            
            # 나머지 컬럼 처리
            for col in df.columns:
                if col not in (key_columns or []):
                    value = self.clean_value(row[col])
                    sentences.append(f"{col} {value}")
            
            chunk = ". ".join(sentences) + ","
            chunks.append(chunk)
        
        # 리스트를 하나의 문자열로 조인
        return " ".join(chunks)
    
    def create_markdown_format(self, df: pd.DataFrame) -> str:
        """
        데이터프레임을 Markdown 형식으로 변환
        """
        markdown_chunks = []
        
        for idx, row in df.iterrows():
            chunk = f"### Record {idx + 1}\n\n"
            
            for col in df.columns:
                value = self.clean_value(row[col])
                chunk += f"- **{col}**: {value}\n"
                # 검색을 위한 별도 태그 추가
                if value:
                    chunk += f"- ##{col.lower().replace(' ', '_')}_{value.lower().replace(' ', '_')}\n"
            
            chunk += "\n---\n"
            markdown_chunks.append(chunk)
        
        return "\n".join(markdown_chunks)

    def create_json_format(self, df: pd.DataFrame) -> str:
        """
        데이터프레임을 JSON 형식으로 변환
        
        Args:
            df: pandas DataFrame
        Returns:
            JSON 형식의 문자열
        """
        try:
            records = []
            
            # DataFrame이 비어있는지 확인
            if df.empty:
                logger.warning("Empty DataFrame received")
                return json.dumps({"records": []}, ensure_ascii=False, indent=2)
            
            for idx, row in df.iterrows():
                try:
                    # 필드 값 변환 시 오류 처리
                    fields = {}
                    for col in df.columns:
                        try:
                            raw_value = row[col]
                            # value가 아닌 raw_value를 사용하여 실제 데이터 값을 가져옴
                            cleaned_value = self.clean_value(raw_value)
                            fields[str(col)] = cleaned_value
                        except Exception as e:
                            logger.warning(f"Error processing column {col}: {str(e)}")
                            fields[str(col)] = ""

                    # 검색 태그 생성
                    search_tags = []
                    for col in df.columns:
                        try:
                            if not pd.isna(row[col]):
                                col_str = str(col).lower().replace(' ', '_')
                                val_str = str(row[col]).lower().replace(' ', '_')
                                if val_str:  # 빈 문자열이 아닌 경우만 태그 추가
                                    search_tags.append(f"{col_str}_{val_str}")
                        except Exception as e:
                            logger.warning(f"Error creating search tag for {col}: {str(e)}")

                    record = {
                        "record_id": int(idx) + 1,
                        "fields": fields,
                        "search_tags": search_tags
                    }
                    records.append(record)
                    
                except Exception as e:
                    logger.error(f"Error processing row {idx}: {str(e)}")
                    continue

            # 최종 JSON 생성
            json_str = json.dumps(
                {"records": records}, 
                ensure_ascii=False, 
                indent=2,
                default=str  # 직렬화할 수 없는 객체 처리
            )
            
            # 결과 검증
            try:
                json.loads(json_str)  # JSON 유효성 검사
                return json_str
            except json.JSONDecodeError as e:
                logger.error(f"Invalid JSON produced: {str(e)}")
                return json.dumps({"records": []}, ensure_ascii=False, indent=2)

        except Exception as e:
            logger.error(f"Error in create_json_format: {str(e)}")
            return json.dumps({"records": []}, ensure_ascii=False, indent=2)    
           
    

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
        return self.clean_text2(text)
    
    def extract_text_from_pptx_pages(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from PowerPoint file pages.

        Args:
            file (Union[str, IO]): Path to the PowerPoint file or a file-like object

        Returns:
            List[Document]: List of Document objects containing extracted text from PowerPoint slides
        """
        documents = []
        temp_file = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file_path)
                prs = Presentation(file_path)
            else:
                # 임시 파일 생성
                temp_file = io.BytesIO(file.read())
                file_name = self.get_file_name(file)
                file_path = file_name  # 파일 객체의 이름을 file_path로 사용
                prs = Presentation(temp_file)

            for slide_number, slide in enumerate(prs.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text += shape.text + "\n"
                
                # 특수 문자 제거
                slide_text = self.clean_text2(slide_text)
                # Document 객체를 생성합니다
                metadata = {
                    "source": os.path.basename(file_name),
                    "file_name": file_name,
                    "page": slide_number
                }
                doc = Document(page_content=slide_text, metadata=metadata)
                documents.append(doc)

            return documents

        except Exception as e:
            logger.debug(f"Error processing PowerPoint file: {str(e)}")
            return []

        finally:
            if temp_file:
                temp_file.close()    
    def chunk_text(self, text, chunk_size):
        return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

    def extract_text_from_file(self, file, file_name):
        """파일에서 텍스트 추출"""
        try:
            if isinstance(file, str):
                file_path = file
                _, ext = os.path.splitext(file_name)
            elif hasattr(file, 'name'):
                file_path = file
                _, ext = os.path.splitext(file_name)
            else:
                raise ValueError("Invalid file input")
            
            ext = ext.lower()
            logger.debug(f"Processing file: {file_name} with extension: {ext}")
            
            # 파일 형식별 처리
            if ext == '.pdf':
                docs = self.extract_text_from_pdf_pages_plumber(file_path, file_name)
            elif ext in ['.docx', '.doc']:
                docs = self.extract_text_from_docx_pages(file_path)
            elif ext in ['.xlsx', '.xls','.csv']:
                docs = self.extract_text_from_xlsx_pages(file_path)
            elif ext in ['.pptx', '.ppt']:
                docs = self.extract_text_from_pptx_pages(file_path)
            elif ext in ['.hwp', '.hwpx']:
                docs = self.hwp_to_pdf(file_path, file_name)
            elif ext in ['.txt']:
                docs = self.extract_text_from_txt(file_path)
            elif ext in ['.md']:
                docs = self.extract_text_from_markdown(file_path)
            elif ext in ['.htm', '.html']:
                docs = self.extract_text_from_html(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
            
            if not docs:
                logger.error(f"No text extracted from file: {file_name}")
                raise ValueError("No text could be extracted from file")
                
            logger.info(f"Successfully extracted text from {file_name}")
            return docs
            
        except Exception as e:
            logger.error(f"Error extracting text from file {file_name}: {str(e)}")
            raise
        
    def hwp_to_pdf(self, hwp_file: Union[str, IO], file_name) -> Optional[List[Document]]:
        """
        Convert HWP file to PDF and extract text. On non-Windows, only extract text.
        
        Args:
            hwp_file (Union[str, IO]): Path to the HWP file or a file-like object
            file_name (str): Name of the file
        Returns:
            Optional[List[Document]]: List of Document objects with extracted text and source
        """
        is_windows = sys.platform.startswith('win')
        if is_windows:            
            return self._hwp_to_pdf_windows(hwp_file, file_name)
        else:            
            return self._extract_text_from_hwp(hwp_file, file_name)

    def _hwp_to_pdf_windows(self, hwp_file: Union[str, IO], file_name) -> Optional[List[Document]]:
        hwp = None
        temp_hwp_path = None
        temp_pdf_path = None
        temp_dir=""
      
        try:
            import pythoncom
            import win32com
            import win32com.client as win32

            pythoncom.CoInitialize()
            hwp = win32.gencache.EnsureDispatch("HWPFrame.HwpObject")
            hwp.RegisterModule("FilePathCheckDLL", "FilePathCheckerModule")

            if isinstance(hwp_file, str):   #파일 객체가 아닌 경우
                hwp_path = hwp_file
                base_name = os.path.splitext(os.path.basename(hwp_path ))[0]                
                temp_dir = os.path.dirname(hwp_path)
        
            else: #파일 객체인 경우
                base_name = os.path.splitext(file_name)[0]                
                temp_dir = os.path.dirname(hwp_file)
                temp_hwp_path = os.path.join(temp_dir, f"{base_name}.hwp")
                with open(temp_hwp_path, 'wb') as temp_file:
                    temp_file.write(hwp_file.read())
                hwp_path = temp_hwp_path

            temp_pdf_path = os.path.join(temp_dir, f"{base_name}.pdf")
            #file_name = f"{base_name}.hwp" #파일 존재 여부를 확인하기 위해 원래 파일명을 셋팅
            
            hwp.Open(hwp_path)
            hwp.SaveAs(temp_pdf_path, "PDF")

            time.sleep(0.5)
            self.handle_completion_dialog()

            time.sleep(0.5)
            self.close_hwp_popups()
            file_name = os.path.basename(temp_pdf_path) # hwp이름을 pdf로 확장자로 변경하여 전달
            #print(f"file path for pdf :{temp_pdf_path}  {file_name}")
            logger.debug(f"PDF file created: {temp_pdf_path}")            
            documents = self.extract_text_from_pdf_pages_plumber(temp_pdf_path, file_name)

            logger.info(f"HWP to PDF conversion completed: {temp_pdf_path}")
            return documents

        except ImportError:
            logger.error("Required Windows libraries not found. Make sure you're on Windows and have the necessary components installed.")
            return None
        except Exception as e:
            logger.error(f"Error during HWP to PDF conversion: {str(e)}")
            return None
        except win32com.client.pywintypes.com_error as e:
            logger.error(f"COM error during HWP to PDF conversion: {str(e)}")
        except IOError as e:
            logger.error(f"IO error during file operations: {str(e)}")
       
        finally:
            if hwp:
                hwp.Quit()
            pythoncom.CoUninitialize()
            
            
            if hwp_path and os.path.exists(hwp_path):
                try:
                    os.remove(hwp_path)
                    logging.info(f"임시 hwp 파일 삭제 완료: {hwp_path}")
                except Exception as e:
                    logging.error(f"임시 hwp 파일 삭제 중 오류 발생: {str(e)}")
           
            if temp_pdf_path and os.path.exists(temp_pdf_path):
                try:
                    os.remove(temp_pdf_path)
                    logging.info(f"임시 PDF 파일 삭제 완료: {temp_pdf_path}")
                except Exception as e:
                    logging.error(f"임시 PDF 파일 삭제 중 오류 발생: {str(e)}")

    def _extract_text_from_hwp(self, hwp_file: Union[str, IO], file_name) -> Optional[List[Document]]:
        try:
            
            if isinstance(hwp_file, str):
                hwp_path = hwp_file
            else:
                temp_dir = tempfile.gettempdir()
                temp_hwp_path = os.path.join(temp_dir, file_name)
                with open(temp_hwp_path, 'wb') as temp_file:
                    temp_file.write(hwp_file.read())
                hwp_path = temp_hwp_path

            hwp = hwp5.HwpFile(hwp_path)
            text = hwp.text().strip()
            
            document = Document(page_content=text, metadata={"source": os.path.basename(file_name)})
            return [document]
        except ImportError:
            logging.error("pyhwp 라이브러리가 설치되지 않았습니다. pip install pyhwp로 설치해주세요.")
            return None
        except Exception as e:
            logging.error(f"HWP 파일 처리 중 오류 발생: {str(e)}")
            return None
        finally:
            if not isinstance(hwp_file, str) and os.path.exists(hwp_path):
                try:
                    os.remove(hwp_path)
                    logging.info(f"임시 HWP 파일 삭제 완료: {hwp_path}")
                except Exception as e:
                    logging.error(f"임시 HWP 파일 삭제 중 오류 발생: {str(e)}")

    # 기존의 handle_completion_dialog, close_hwp_popups, extract_text_from_pdf_pages 메서드는 그대로 유지

# 사용 예:
# converter = HwpConverter()
# documents = converter.hwp_to_pdf("path/to/file.hwp", "file.hwp")

    def close_hwp_popups(self):
        def callback(hwnd, _):
            if win32gui.GetWindowText(hwnd) == "한글":
                win32gui.PostMessage(hwnd, win32con.WM_CLOSE, 0, 0)
        
        win32gui.EnumWindows(callback, None)
    
    def handle_completion_dialog(self):
        def callback(hwnd, _):
            if "완료" in win32gui.GetWindowText(hwnd):
                # '확인' 버튼 클릭
                ok_button = win32gui.FindWindowEx(hwnd, 0, "Button", "확인")
                win32gui.PostMessage(ok_button, win32con.WM_LBUTTONDOWN, 0, 0)
                win32gui.PostMessage(ok_button, win32con.WM_LBUTTONUP, 0, 0)
        
        win32gui.EnumWindows(callback, None)
        

    def extract_text_from_pdf(self, file):
        text = ""
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")
            # 파일 객체가 읽기 모드로 열려 있는지 확인
        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")
    
        # 파일 경로인 경우 파일 열기
        if isinstance(file, str):
            file = open(file, 'rb')
        
        try:
            reader = PyPDF2.PdfReader(file)
            documents = []
            for page_num, page in enumerate(reader.pages,1):
                text += page.extract_text()
                text = self.clean_text2(text)
                metadata = {
                    "source": os.path.basename(file.name if hasattr(file, 'name') else 'unknown'),
                    "file_name": getattr(file, 'name', 'unknown'),
                    "page": page_num
                }
                doc = Document(page_content=text, metadata=metadata)
                documents.append(doc)
        finally:
            # 우리가 파일을 열었다면 닫아줍니다
            if isinstance(file, io.IOBase) and not isinstance(file, io.BytesIO):
                file.close()        
        return documents
    
    def get_file_name(self, file_object):
        name = getattr(file_object, 'name', None)
        if name:
            return os.path.basename(name)
        return "Unknown"
    
    def extract_text_from_markdown(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from Markdown file.

        Args:
            file (Union[str, IO]): Path to the Markdown file or a file-like object

        Returns:
            List[Document]: List containing a single Document object with extracted text
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        file_obj = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file)
                file_obj = open(file, 'r', encoding='utf-8')
            else:
                file_obj = file
                file_name = self.get_file_name(file)
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우

            content = file_obj.read()
            html = markdown.markdown(content)
            text = self.extract_from_html(html)

            metadata = {
                "source": os.path.basename(file_name),
                "file_name": file_name,
                "file_type": "markdown"
            }
            doc = Document(page_content=text, metadata=metadata)
            documents.append(doc)

            return documents

        except Exception as e:
            logger.debug(f"Error processing Markdown file: {str(e)}")
            return []

        finally:
            if file_obj and isinstance(file_obj, io.IOBase) and not isinstance(file_obj, io.StringIO):
                file_obj.close()

    def extract_text_from_html(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from HTML file.

        Args:
            file (Union[str, IO]): Path to the HTML file or a file-like object

        Returns:
            List[Document]: List containing a single Document object with extracted text
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        file_obj = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file)
                file_obj = open(file, 'r', encoding='utf-8')
            else:
                file_obj = file
                file_name = self.get_file_name(file)
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우

            content = file_obj.read()
            text = self.extract_from_html(content)

            metadata = {
                "source": os.path.basename(file_name),
                "file_name": file_name,
                "file_type": "html"
            }
            doc = Document(page_content=text, metadata=metadata)
            documents.append(doc)

            return documents

        except Exception as e:
            logger.debug(f"Error processing HTML file: {str(e)}")
            return []

        finally:
            if file_obj and isinstance(file_obj, io.IOBase) and not isinstance(file_obj, io.StringIO):
                file_obj.close()

    def extract_text_from_txt(self, file: Union[str, IO]) -> List[Document]:
        """
        Extract text from TXT file.

        Args:
            file (Union[str, IO]): Path to the TXT file or a file-like object

        Returns:
            List[Document]: List containing a single Document object with extracted text
        """
        if not isinstance(file, (str, io.IOBase)):
            raise TypeError("file must be a string path or a file-like object")

        if isinstance(file, io.IOBase) and not file.readable():
            raise ValueError("file must be opened in read mode")

        documents = []
        file_obj = None

        try:
            if isinstance(file, str):
                file_path = file
                file_name = os.path.basename(file)
                file_obj = open(file, 'r', encoding='utf-8')
            else:
                file_obj = file
                file_name = self.get_file_name(file)
                file_path = file_name  # Streamlit의 file_uploader를 사용할 경우

            content = file_obj.read()
            text = content.strip()

            metadata = {
                "source": os.path.basename(file_name),
                "file_name": file_name,
                "file_type": "txt"
            }
            doc = Document(page_content=text, metadata=metadata)
            documents.append(doc)

            return documents

        except Exception as e:
            logger.debug(f"Error processing TXT file: {str(e)}")
            return []

        finally:
            if file_obj and isinstance(file_obj, io.IOBase) and not isinstance(file_obj, io.StringIO):
                file_obj.close()

    def extract_from_html(self, content: str) -> str:
        """Extract text from HTML content."""
        soup = BeautifulSoup(content, 'html.parser')
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text()
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        return text